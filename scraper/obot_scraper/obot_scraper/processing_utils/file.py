import re

import pypdf
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook

from langchain_text_splitters import RecursiveCharacterTextSplitter


def extract_text_from_pdf(file):
    try:
        reader = pypdf.PdfReader(file)  # Create a PdfReader object
    except Exception as e:
        raise ValueError("Failed to extract text from PDF. This may not be a pdf") from e

    text = ""
    for page in range(len(reader.pages)):
        text += reader.pages[page].extract_text()  # Extract text from each page
    
    return text
    

def extract_text_from_docx(file):
    doc = Document(file)  # Create a Document object

    text = ""
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"  # Extract text from each paragraph
    
    return text

def extract_text_from_pptx(file):
    prs = Presentation(file)  # Create a Presentation object

    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"  # Extract text from each shape
    
    return text

def extract_text_from_xlsx(file):
    wb = load_workbook(file)  # Create a Workbook object
    ws = wb.active  # Get the active worksheet

    text = ""
    for row in ws.iter_rows(values_only=True):
        text += ' | '.join([str(cell) for cell in row if cell is not None]) + "\n"  # Extract text from each row
    
    return text


def file_to_chunks(file, url):
    if url.endswith('.pdf') or url.startswith(("https://drive.google.com", "https://drive.usercontent.google.com")):
        text = extract_text_from_pdf(file)
    elif url.endswith('.docx'):
        text = extract_text_from_docx(file)
    elif url.endswith('.pptx'):
        text = extract_text_from_pptx(file)
    elif url.endswith('.xlsx'):
        text = extract_text_from_xlsx(file)
    else:
        raise ValueError("Unsupported file type")

    # More than two newlines to two new lines
    text = re.sub(r'\n{3,}', '\n\n', text)
    
    # Split the text into chunks
    splitter = RecursiveCharacterTextSplitter(
        chunk_size = 3000,
        chunk_overlap = 0,
        length_function = lambda x: len(x.split()),
        separators = ["\n\n", "\n", ".", " ", ""]
    )

    return splitter.split_text(text)